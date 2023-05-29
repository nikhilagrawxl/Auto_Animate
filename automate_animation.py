import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches

# Open the PowerPoint presentation
presentation = Presentation('presentation.pptx')

# Access the first slide
slide = presentation.slides[0]  # Modify slide index as needed

# Create a new shape and apply an animation effect
left = Inches(1)
top = Inches(1)
width = Inches(11)
height = Inches(5)
shape = slide.shapes.add_shape(
    1,  # Rectangle shape
    left=left,
    top=top,
    width=width,
    height=height
)

# Add a text box and adjust its position
left = Inches(2)
top = Inches(2)
width = Inches(7)
height = Inches(1)
text_box = slide.shapes.add_textbox(
    left=left,
    top=top,
    width=width,
    height=height
)
text_frame = text_box.text_frame
text_frame.text = 'Hello students Today we will be solving a question on chapter Moving Charges and Magnetism and the question goes like A three point zero centimeter wire carrying a current of ten amperes is placed inside a solenoid perpendicular to its axis. The magnetic field inside the solenoid is given to be zero point two seven teslas What is the magnetic force on the wire'

# Save the modified presentation
presentation.save('modified_presentation.pptx')
